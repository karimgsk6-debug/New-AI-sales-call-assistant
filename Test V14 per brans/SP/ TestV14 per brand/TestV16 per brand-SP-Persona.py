import streamlit as st
from PIL import Image
import requests
from io import BytesIO, BytesIO as io_bytes
import groq
from groq import Groq
import streamlit.components.v1 as components
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Initialize Groq client ---
client = Groq(api_key="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")

# --- Initialize session state ---
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# --- Language selector ---
language = st.radio("Select Language / اختر اللغة", options=["English", "العربية"])

# --- GSK logo ---
logo_local_path = "images/gsk_logo.png"
logo_fallback_url = "https://www.tungsten-network.com/wp-content/uploads/2020/05/GSK_Logo_Full_Colour_RGB.png"

col1, col2 = st.columns([1, 5])
with col1:
    try:
        logo_img = Image.open(logo_local_path)
        st.image(logo_img, width=120)
    except Exception:
        st.image(logo_fallback_url, width=120)
with col2:
    st.title("🧠 AI Sales Call Assistant")

# --- GSK brand mappings ---
gsk_brands = {
    "Trelegy": "https://example.com/trelegy-leaflet",
    "Shingrix": "https://example.com/shingrix-leaflet",
    "Zejula": "https://example.com/zejula-leaflet",
}

gsk_brands_images = {
    "Trelegy": "https://www.example.com/trelegy.png",
    "Shingrix": "https://www.oma-apteekki.fi/WebRoot/NA/Shops/na/67D6/48DA/D0B0/D959/ECAF/0A3C/0E02/D573/3ad67c4e-e1fb-4476-a8a0-873423d8db42_3Dimage.png",
    "Zejula": "https://cdn.salla.sa/QeZox/eyy7B0bg8D7a0Wwcov6UshWFc04R6H8qIgbfFq8u.png",
}

# --- Filters ---
race_segments = [
    "R – Reach: Did not start to prescribe yet and Don't believe that vaccination is his responsibility.",
    "A – Acquisition: Prescribe to patient who initiate discussion about the vaccine but Convinced about Shingrix data.",
    "C – Conversion: Proactively initiate discussion with specific patient profile but For other patient profiles he is not prescribing yet.",
    "E – Engagement: Proactively prescribe to different patient profiles"
]

doctor_barriers = [
    "1 - HCP does not consider HZ as risk for the selected patient profile",
    "2 - HCP thinks there is no time to discuss preventive measures with the patients",
    "3 - HCP thinks about cost considerations",
    "4 - HCP is not convinced that HZ Vx is effective in reducing the burden",
    "5 - Accessibility (POVs)"
]

objectives = ["Awareness", "Adoption", "Retention"]
specialties = ["General Practitioner", "Cardiologist", "Dermatologist", "Endocrinologist", "Pulmonologist"]
personas = [
    "Uncommitted Vaccinator – Not engaged, poor knowledge, least likely to prescribe vaccines (26%)",
    "Reluctant Efficiency – Do not see vaccinating 50+ as part of role, least likely to believe in impact (12%)",
    "Patient Influenced – Aware of benefits but prescribes only if patient requests (26%)",
    "Committed Vaccinator – Very positive, motivated, prioritizes vaccination & sets example (36%)"
]
gsk_approaches = [
    "Use data-driven evidence",
    "Focus on patient outcomes",
    "Leverage storytelling techniques",
]
sales_call_flow = ["Prepare", "Engage", "Create Opportunities", "Influence", "Drive Impact", "Post Call Analysis"]

# --- Sidebar Filters ---
st.sidebar.header("Filters & Options")
brand = st.sidebar.selectbox("Select Brand / اختر العلامة التجارية", options=list(gsk_brands.keys()))
segment = st.sidebar.selectbox("Select RACE Segment / اختر شريحة RACE", race_segments)
barrier = st.sidebar.multiselect("Select Doctor Barrier / اختر حاجز الطبيب", options=doctor_barriers, default=[])
objective = st.sidebar.selectbox("Select Objective / اختر الهدف", objectives)
specialty = st.sidebar.selectbox("Select Doctor Specialty / اختر تخصص الطبيب", specialties)
persona = st.sidebar.selectbox("Select HCP Persona / اختر شخصية الطبيب", personas)
response_length_options = ["Short", "Medium", "Long"]
response_tone_options = ["Formal", "Casual", "Friendly", "Persuasive"]
response_length = st.sidebar.selectbox("Select Response Length / اختر طول الرد", response_length_options)
response_tone = st.sidebar.selectbox("Select Response Tone / اختر نبرة الرد", response_tone_options)
interface_mode = st.sidebar.radio("Interface Mode / اختر واجهة", ["Chatbot", "Card Dashboard", "Flow Visualization"])

# --- Load brand image safely ---
image_path = gsk_brands_images.get(brand)
try:
    if image_path.startswith("http"):
        response = requests.get(image_path)
        img = Image.open(BytesIO(response.content))
    else:
        img = Image.open(image_path)
    st.image(img, width=200)
except Exception:
    st.warning(f"⚠️ Could not load image for {brand}. Using placeholder.")
    st.image("https://via.placeholder.com/200x100.png?text=No+Image", width=200)

# --- Clear chat button ---
if st.button("🗑️ Clear Chat / مسح المحادثة"):
    st.session_state.chat_history = []

# --- Chat container ---
chat_container = st.container()
placeholder_text = "Type your message..." if language == "English" else "اكتب رسالتك..."
user_input = st.text_area(placeholder_text, key="user_input", height=80)

# --- Send button ---
if st.button("🚀 Send / أرسل") and user_input.strip():
    with st.spinner("Generating AI response... / جارٍ إنشاء الرد"):
        st.session_state.chat_history.append({"role": "user", "content": user_input})

        approaches_str = "\n".join(gsk_approaches)
        flow_str = " → ".join(sales_call_flow)

        prompt = f"""
Language: {language}
You are an expert GSK sales assistant. 
User input: {user_input}
RACE Segment: {segment}
Doctor Barrier: {', '.join(barrier) if barrier else 'None'}
Objective: {objective}
Brand: {brand}
Doctor Specialty: {specialty}
HCP Persona: {persona}
Approved GSK Sales Approaches:
{approaches_str}
Sales Call Flow Steps:
{flow_str}
Use APACT (Acknowledge → Probing → Answer → Confirm → Transition) technique for handling objections.
Response Length: {response_length}
Response Tone: {response_tone}
Provide actionable suggestions tailored to this persona, following the selected length and tone, in a friendly and professional manner.
"""

        # Call Groq API
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[
                {"role": "system", "content": f"You are a helpful sales assistant chatbot that responds in {language}."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7
        )

        ai_output = response.choices[0].message.content
        st.session_state.chat_history.append({"role": "ai", "content": ai_output})

        # --- Download options ---
        if ai_output:
            st.markdown("### 📥 Download AI Response")
            
            # Word
            doc = Document()
            doc.add_heading("AI Sales Call Response", 0)
            doc.add_paragraph(ai_output)
            word_buffer = io_bytes()
            doc.save(word_buffer)
            st.download_button("Download as Word (.docx)", word_buffer.getvalue(), file_name="AI_Response.docx")

            # PowerPoint
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = "AI Sales Call Response"
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
            content_frame = content_box.text_frame
            content_frame.text = ai_output
            ppt_buffer = io_bytes()
            prs.save(ppt_buffer)
            st.download_button("Download as PowerPoint (.pptx)", ppt_buffer.getvalue(), file_name="AI_Response.pptx")

# --- Display chat history ---
with chat_container:
    if interface_mode == "Chatbot":
        st.subheader("💬 Chatbot Interface")
        for msg in st.session_state.chat_history:
            if msg["role"] == "user":
                st.markdown(f"<div style='text-align:right; background:#d1e7dd; padding:10px; border-radius:12px; margin:10px 0;'>{msg['content']}</div>", unsafe_allow_html=True)
            else:
                st.markdown(f"<div style='text-align:left; background:#f0f2f6; padding:15px; border-radius:12px; margin:10px 0; box-shadow:2px 2px 5px rgba(0,0,0,0.1);'>{msg['content']}</div>", unsafe_allow_html=True)

    elif interface_mode == "Card Dashboard":
        st.subheader("📊 Card-Based Dashboard")
        segments_list = ["Evidence-Seeker", "Skeptic", "Time-Pressured"]
        for seg in segments_list:
            with st.expander(f"{seg} Segment"):
                st.write(f"Suggested approach for {seg} with {', '.join(barrier) if barrier else 'None'} barriers selected.")
                st.progress(70)
                st.button(f"Next Suggestion for {seg}")

    elif interface_mode == "Flow Visualization":
        st.subheader("🔗 HCP Engagement Flow")
        html_content = f"""
        <div style='font-family:sans-serif; background:#f0f2f6; padding:20px; border-radius:10px;'>
            <h3>{persona} Segment</h3>
            <p><b>Behavior:</b> {', '.join(barrier) if barrier else 'None'}</p>
            <p><b>Brand:</b> {brand}</p>
            <p><b>Sales Flow:</b> {flow_str}</p>
            <p><b>Tone:</b> {response_tone}</p>
            <p><b>AI Suggestion:</b> Example probing question or APACT objection handling approach here...</p>
        </div>
        """
        components.html(html_content, height=300)

# --- Brand leaflet ---
st.markdown(f"[Brand Leaflet - {brand}]({gsk_brands[brand]})")
